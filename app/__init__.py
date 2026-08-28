"""
WARP Tool - Flask Application

Routes:
  GET  /                               - Redirect zum letzten Projekt oder Neuanlage
  POST /project/new                    - Neues Projekt anlegen
  GET  /project/<id>                   - Fragenkatalog für ein Projekt
  POST /project/<id>/answer            - AJAX: Antwort/Notiz speichern
  POST /project/<id>/info              - AJAX: Projektinfos speichern
  POST /project/<id>/delete            - Projekt löschen
  POST /project/<id>/report/html       - HTML-Vorschau
  POST /project/<id>/report            - PDF-Download
  GET  /login                          - Login
  POST /login                          - Login verarbeiten
  GET  /register                       - Registrierung
  POST /register                       - Registrierung verarbeiten
  GET  /logout                         - Logout
  GET  /admin                          - Admin-Übersicht
  GET  /admin/questions                - Fragenkatalog verwalten
  POST /admin/question/new             - Neue Frage anlegen
  POST /admin/question/<qid>/edit      - Frage bearbeiten
  POST /admin/question/<qid>/delete    - Frage löschen
  POST /admin/category/new             - Neue Kategorie anlegen
  POST /admin/category/<cid>/edit      - Kategorie bearbeiten
  POST /admin/category/<cid>/delete    - Kategorie löschen
  POST /api/inbox                      - Webhook: IMPULSE-Nachricht empfangen (API-Key)
  GET  /api/inbox/count                - AJAX: Anzahl neuer Nachrichten
  GET  /inbox                          - Admin-Postkorb
  POST /inbox/<mid>/claim              - Nachricht übernehmen
  POST /inbox/<mid>/done               - Nachricht als erledigt markieren
  POST /inbox/<mid>/release            - Nachricht freigeben
  POST /inbox/<mid>/delete             - Nachricht löschen
"""

from __future__ import annotations

import io
import json as _json
import math
import os
import uuid
import datetime as dt
from pathlib import Path
from collections import OrderedDict
from typing import Any, Dict, List

from flask import (
    Flask, render_template, request, send_file, abort,
    redirect, url_for, jsonify, flash,
)
from flask_login import (
    LoginManager, login_required, login_user, logout_user, current_user,
)
from flask_wtf.csrf import CSRFProtect
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

from .models import db, User, Project, Answer, Category, Question, InboxMessage, GeneratedDocument, WorkshopQuestion
from .data.questions import ANSWER_OPTIONS, WARP_LEVEL_ORDER, CATEGORIES as _PY_CATS

# Lookup: question_id → recommendations dict (low/mid), built from questions.py
_REC_LOOKUP: dict = {
    q["id"]: q["recommendations"]
    for cat in _PY_CATS
    for q in cat["questions"]
    if q.get("recommendations")
}

ROOT = Path(__file__).parent.parent


def _donut_segments(
    segments: List[tuple],
    outer_r: float = 42.0,
    inner_r: float = 27.0,
    cx: float = 50.0,
    cy: float = 50.0,
) -> List[Dict[str, Any]]:
    """Bereitet Segmente für ein SVG-Donut-Diagramm als explizite Ring-Pfade auf.

    segments: Liste von (id, label, color, count)-Tupeln.
    Jedes Segment bekommt ein fertiges SVG-<path>-'d'-Attribut (Ringsegment
    zwischen innerem und äußerem Radius), berechnet über exakte Winkel im
    Uhrzeigersinn ab 12 Uhr — bewusst über direkte Pfadgeometrie statt der
    fehleranfälligen stroke-dasharray/rotate-Technik auf <circle>-Elementen.
    """
    def point(angle_deg: float, r: float) -> tuple:
        rad = math.radians(angle_deg)
        return (cx + r * math.sin(rad), cy - r * math.cos(rad))

    total = sum(s[3] for s in segments) or 1
    out: List[Dict[str, Any]] = []
    start_deg = 0.0
    for sid, label, color, count in segments:
        pct = count / total * 100
        sweep_deg = min(pct / 100 * 360, 359.99)  # 360° entartet den Arc-Befehl
        end_deg = start_deg + sweep_deg
        path = ""
        if sweep_deg > 0.01:
            large_arc = 1 if sweep_deg > 180 else 0
            ox1, oy1 = point(start_deg, outer_r)
            ox2, oy2 = point(end_deg, outer_r)
            ix2, iy2 = point(end_deg, inner_r)
            ix1, iy1 = point(start_deg, inner_r)
            path = (
                f"M {ox1:.3f},{oy1:.3f} "
                f"A {outer_r:.3f},{outer_r:.3f} 0 {large_arc} 1 {ox2:.3f},{oy2:.3f} "
                f"L {ix2:.3f},{iy2:.3f} "
                f"A {inner_r:.3f},{inner_r:.3f} 0 {large_arc} 0 {ix1:.3f},{iy1:.3f} Z"
            )
        out.append({
            "id": sid,
            "label": label,
            "color": color,
            "pct": round(pct, 1),
            "count": count,
            "path": path,
        })
        start_deg = end_deg
    return out

csrf = CSRFProtect()
limiter = Limiter(key_func=get_remote_address, default_limits=[])


# ---------------------------------------------------------------------------
# DB-Seed: Fragen aus Python-Code beim ersten Start laden
# ---------------------------------------------------------------------------

def _seed_questions_if_needed() -> None:
    from .data.questions import CATEGORIES as PY_CATS
    count = db.session.execute(
        db.select(db.func.count()).select_from(Category)
    ).scalar()
    if count and count > 0:
        return
    for sort_idx, cat_data in enumerate(PY_CATS):
        cat = Category(
            id=cat_data["id"],
            title=cat_data["title"],
            parent=cat_data["parent"],
            description=cat_data.get("description", ""),
            sort_order=sort_idx,
        )
        db.session.add(cat)
        for q_idx, q_data in enumerate(cat_data["questions"]):
            q = Question(
                id=q_data["id"],
                category_id=cat_data["id"],
                text=q_data["text"],
                hint=q_data.get("hint"),
                is_new=q_data.get("new", False),
                sort_order=q_idx,
            )
            db.session.add(q)
    db.session.commit()
    q_total = sum(len(c["questions"]) for c in PY_CATS)
    print(f"[WARP] Fragenkatalog ({q_total} Fragen, {len(PY_CATS)} Kategorien) in DB geladen.")


def create_app() -> Flask:
    app = Flask(__name__, template_folder="templates", static_folder="static")

    app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", "warp-dev-secret-key-change-in-production")

    db_url = os.environ.get("DATABASE_URL", f"sqlite:///{ROOT / 'warp.db'}")
    if db_url.startswith("postgres://"):
        db_url = db_url.replace("postgres://", "postgresql://", 1)
    if db_url.startswith("postgresql://") and "sslmode" not in db_url:
        db_url += "?sslmode=require"
    app.config["SQLALCHEMY_DATABASE_URI"] = db_url
    app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False

    db.init_app(app)

    # Session & Cookie Security
    from datetime import timedelta
    is_production = bool(os.environ.get("RENDER"))
    app.config.setdefault("WTF_CSRF_ENABLED", True)
    app.config["SESSION_COOKIE_HTTPONLY"] = True
    app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
    app.config["SESSION_COOKIE_SECURE"] = is_production
    app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(hours=8)
    app.config["RATELIMIT_ENABLED"] = os.environ.get("RATELIMIT_ENABLED", "true").lower() != "false"
    app.config.setdefault("RATELIMIT_STORAGE_URI", "memory://")

    csrf.init_app(app)
    limiter.init_app(app)

    login_manager = LoginManager()
    login_manager.init_app(app)
    login_manager.login_view = "login"
    login_manager.login_message = "Bitte melden Sie sich an."

    @login_manager.user_loader
    def load_user(user_id: str):
        return db.session.get(User, int(user_id))

    with app.app_context():
        db.create_all()

        try:
            admin_user = db.session.execute(
                db.select(User).where(User.username == "admin")
            ).scalar_one_or_none()
            if not admin_user:
                default = User(username="admin", display_name="Administrator", role='superuser')
                default.set_password("warp2024")
                db.session.add(default)
                db.session.commit()
                print("[WARP] Standard-Admin erstellt: admin / warp2024")
            elif not admin_user.is_superuser:
                admin_user.role = 'superuser'
                db.session.commit()
        except Exception as e:
            db.session.rollback()
            print(f"[WARP] Admin-Init übersprungen: {e}")

        try:
            _seed_questions_if_needed()
        except Exception as e:
            db.session.rollback()
            print(f"[WARP] Fragen-Seed übersprungen: {e}")

    score_lookup = {opt["id"]: opt["score"] for opt in ANSWER_OPTIONS}
    label_lookup = {opt["id"]: opt["label"] for opt in ANSWER_OPTIONS}

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _get_project_or_403(pid: int) -> Project:
        project = db.session.get(Project, pid)
        if not project:
            abort(404)
        if not current_user.is_admin and project.user_id != current_user.id:
            abort(403)
        return project

    @app.errorhandler(403)
    def handle_403(e):
        return render_template("403.html"), 403

    def _db_question_count(catalog_type: str = 'assessment') -> int:
        return db.session.execute(
            db.select(db.func.count(Question.id))
            .join(Category, Question.category_id == Category.id)
            .where(Category.catalog_type == catalog_type)
        ).scalar() or 0

    def _load_categories_from_db(catalog_type: str = 'assessment') -> List[Dict]:
        cats = db.session.execute(
            db.select(Category)
            .where(Category.catalog_type == catalog_type)
            .order_by(Category.sort_order)
        ).scalars().all()
        return [
            {
                "id": c.id,
                "title": c.title,
                "parent": c.parent,
                "description": c.description,
                "questions": [
                    {"id": q.id, "text": q.text, "hint": q.hint, "new": q.is_new,
                     "recommendations": _REC_LOOKUP.get(q.id, {})}
                    for q in c.questions
                ],
            }
            for c in cats
        ]

    def _upsert_answer(project_id: int, question_id: str,
                       answer_id: str | None, note: str | None) -> None:
        row = db.session.execute(
            db.select(Answer).where(
                Answer.project_id == project_id,
                Answer.question_id == question_id,
            )
        ).scalar_one_or_none()
        if row:
            if answer_id is not None:
                row.answer_id = answer_id or None
            if note is not None:
                row.note = note.strip() or None
        else:
            row = Answer(
                project_id=project_id,
                question_id=question_id,
                answer_id=answer_id or None,
                note=(note.strip() if note else None),
            )
            db.session.add(row)
        db.session.commit()

    def _collect_answers(form) -> Dict[str, str]:
        out: Dict[str, str] = {}
        for key, value in form.items():
            if key.startswith("answer-") and value:
                qid = key[len("answer-"):]
                if value in score_lookup:
                    out[qid] = value
        return out

    def _collect_notes(form) -> Dict[str, str]:
        out: Dict[str, str] = {}
        for key, value in form.items():
            if key.startswith("note-"):
                qid = key[len("note-"):]
                if value.strip():
                    out[qid] = value.strip()
        return out

    def _build_report_context(form, customer: bool = False) -> Dict[str, Any]:
        LEVEL_THRESHOLD = 70

        LEVEL_NUMBERS = {
            "Stufe 2 - Managed": 2,
            "Stufe 3 - Defined": 3,
            "Stufe 4 - Measured": 4,
            "Stufe 5 - Optimization": 5,
        }

        answers = _collect_answers(form)
        notes = _collect_notes(form)
        project_name = form.get("project_name", "").strip() or "Unbenannte Analyse"
        project_owner = form.get("project_owner", "").strip()
        project_date = form.get("project_date", "").strip() or dt.date.today().strftime("%d.%m.%Y")

        db_categories = _load_categories_from_db()

        level_total_map: Dict[str, int] = {
            lv: sum(len(cat["questions"]) for cat in db_categories if cat["parent"] == lv)
            for lv in WARP_LEVEL_ORDER
        }
        total_questions = _db_question_count()

        categories: List[Dict[str, Any]] = []
        all_scores_sum: float = 0.0
        level_scores_sum: Dict[str, float] = {lv: 0.0 for lv in WARP_LEVEL_ORDER}
        overall_distribution: "OrderedDict[str, int]" = OrderedDict(
            (opt["id"], 0) for opt in ANSWER_OPTIONS
        )

        for cat in db_categories:
            cat_scores_sum: float = 0.0
            cat_questions: List[Dict[str, Any]] = []
            answered = 0
            distribution: "OrderedDict[str, int]" = OrderedDict(
                (opt["id"], 0) for opt in ANSWER_OPTIONS
            )

            for q in cat["questions"]:
                aid = answers.get(q["id"])
                if aid:
                    answered += 1
                    distribution[aid] += 1
                    overall_distribution[aid] += 1
                    score = score_lookup[aid]
                    cat_scores_sum += score
                    all_scores_sum += score
                    level_scores_sum[cat["parent"]] += score

                cat_questions.append({
                    "id": q["id"],
                    "text": q["text"],
                    "hint": q.get("hint"),
                    "answer_id": aid,
                    "answer_label": label_lookup.get(aid),
                    "score": score_lookup.get(aid),
                    "note": notes.get(q["id"], ""),
                })

            cat_total = len(cat["questions"])
            cat_avg = cat_scores_sum / cat_total if cat_total > 0 else 0.0
            categories.append({
                "id": cat["id"],
                "title": cat["title"],
                "description": cat["description"],
                "parent": cat["parent"],
                "questions": cat_questions,
                "answered": answered,
                "total": cat_total,
                "score": round(cat_avg, 1),
                "distribution": [
                    {"id": opt["id"], "label": opt["label"], "count": distribution[opt["id"]]}
                    for opt in ANSWER_OPTIONS
                ],
            })

        overall = round(all_scores_sum / total_questions, 1) if total_questions > 0 else 0.0
        answered_count = sum(c["answered"] for c in categories)

        warp_levels: List[Dict[str, Any]] = [{
            "number": 1, "name": "Initial", "full_name": "Stufe 1 – Initial",
            "score": None, "categories": [], "achieved": True,
        }]
        all_lower_achieved = True
        for level_name in WARP_LEVEL_ORDER:
            lv_total = level_total_map.get(level_name, 0)
            lv_score = round(level_scores_sum[level_name] / lv_total, 1) if lv_total > 0 else 0.0
            lv_cats = [c for c in categories if c["parent"] == level_name]
            num = LEVEL_NUMBERS[level_name]
            short_name = level_name.split(" - ")[1]
            achieved = all_lower_achieved and (lv_score >= LEVEL_THRESHOLD)
            if not achieved:
                all_lower_achieved = False
            warp_levels.append({
                "number": num, "name": short_name,
                "full_name": level_name.replace(" - ", " – "),
                "score": lv_score, "categories": lv_cats, "achieved": achieved,
            })

        warp_level_reached = max(ld["number"] for ld in warp_levels if ld["achieved"])
        warp_level_label = next(ld["full_name"] for ld in warp_levels if ld["number"] == warp_level_reached)

        # Nächste, noch nicht erreichte Stufe — Grundlage für den Kundenreport
        # ("Ausblick & Handlungsempfehlungen: wie erreiche ich die nächste Stufe?").
        next_level = next(
            (ld for ld in warp_levels if ld["number"] > 1 and not ld["achieved"]),
            None,
        )
        customer_recs: List[Dict[str, Any]] = []
        next_level_partial_count = 0
        if next_level:
            for cat in next_level["categories"]:
                entries = [
                    {"text": q["text"], "recommendation": _REC_LOOKUP.get(q["id"], {}).get("low")}
                    for q in cat["questions"]
                    if q["answer_id"] == "nicht" and _REC_LOOKUP.get(q["id"], {}).get("low")
                ]
                if entries:
                    gap_summary = "; ".join(e["text"].rstrip(".") for e in entries) + "."
                    body = " ".join(e["recommendation"] for e in entries)
                    customer_recs.append({
                        "title": cat["title"],
                        "count": len(entries),
                        "gap_summary": gap_summary,
                        "body": body,
                    })
                next_level_partial_count += sum(
                    1 for q in cat["questions"] if q["answer_id"] in ("kaum", "teil")
                )

        # KPI-Kuchendiagramme für den Kundenreport (nur dort benötigt).
        distribution_pie: List[Dict[str, Any]] | None = None
        next_level_pie: List[Dict[str, Any]] | None = None
        if customer:
            dist_colors = {"voll": "#14B26A", "teil": "#5BC5DD", "kaum": "#F2A73B", "nicht": "#E84A4A"}
            dist_segments = [
                (opt["id"], opt["label"], dist_colors[opt["id"]], overall_distribution[opt["id"]])
                for opt in ANSWER_OPTIONS
            ]
            unanswered = total_questions - answered_count
            if unanswered > 0:
                dist_segments.append(("offen", "Nicht beantwortet", "#E6E6E6", unanswered))
            distribution_pie = _donut_segments(dist_segments)

            if next_level:
                achieved_pct = next_level["score"]
                gap_pct = max(LEVEL_THRESHOLD - achieved_pct, 0)
                headroom_pct = max(100 - LEVEL_THRESHOLD, 0)
                next_level_pie = _donut_segments([
                    ("erreicht", "Erreicht", "#451DC7", achieved_pct),
                    ("luecke", "Lücke zum Zielwert", "#F2A73B", gap_pct),
                    ("rest", "Darüber hinaus", "#E6E6E6", headroom_pct),
                ])

        return {
            "project_name": project_name,
            "project_owner": project_owner,
            "project_date": project_date,
            "generated_at": dt.datetime.now().strftime("%d.%m.%Y, %H:%M"),
            "categories": categories,
            "overall_score": overall,
            "answered_count": answered_count,
            "total_count": total_questions,
            "maturity_label": warp_level_label,
            "maturity_index": warp_level_reached,
            "warp_levels": warp_levels,
            "warp_level_reached": warp_level_reached,
            "level_threshold": LEVEL_THRESHOLD,
            "answer_options": ANSWER_OPTIONS,
            "customer": customer,
            "next_level": next_level,
            "customer_recs": customer_recs,
            "customer_recs_count": sum(c["count"] for c in customer_recs),
            "distribution_pie": distribution_pie,
            "next_level_pie": next_level_pie,
            "next_level_partial_count": next_level_partial_count,
        }

    # ------------------------------------------------------------------
    # Auth Routes
    # ------------------------------------------------------------------

    @app.route("/register", methods=["GET", "POST"])
    def register():
        if current_user.is_authenticated:
            return redirect(url_for("index"))
        error = None
        if request.method == "POST":
            username = request.form.get("username", "").strip()
            display_name = request.form.get("display_name", "").strip()
            password = request.form.get("password", "")
            password2 = request.form.get("password2", "")
            if not username:
                error = "Benutzername darf nicht leer sein."
            elif len(password) < 6:
                error = "Passwort muss mindestens 6 Zeichen haben."
            elif password != password2:
                error = "Passwörter stimmen nicht überein."
            else:
                existing = db.session.execute(
                    db.select(User).where(User.username == username)
                ).scalar_one_or_none()
                if existing:
                    error = f"Benutzername '{username}' ist bereits vergeben."
                else:
                    u = User(username=username, display_name=display_name or None)
                    u.set_password(password)
                    db.session.add(u)
                    db.session.commit()
                    login_user(u)
                    return redirect(url_for("index"))
        return render_template("register.html", error=error)

    @app.route("/login", methods=["GET", "POST"])
    @limiter.limit("20 per minute")
    def login():
        if current_user.is_authenticated:
            return redirect(url_for("index"))
        error = None
        if request.method == "POST":
            username = request.form.get("username", "").strip()
            password = request.form.get("password", "")
            user = db.session.execute(
                db.select(User).where(User.username == username)
            ).scalar_one_or_none()
            if user and user.check_password(password):
                if user.is_locked:
                    error = "Ihr Konto wurde gesperrt. Bitte wenden Sie sich an den Administrator."
                else:
                    login_user(user)
                    next_page = request.args.get("next", "")
                    if not next_page.startswith("/"):
                        next_page = url_for("index")
                    return redirect(next_page)
            elif not error:
                error = "Ungültiger Benutzername oder Passwort."
        return render_template("login.html", error=error)

    @app.route("/logout")
    @login_required
    def logout():
        logout_user()
        return redirect(url_for("login"))

    # ------------------------------------------------------------------
    # Project Routes
    # ------------------------------------------------------------------

    @app.route("/")
    @login_required
    def index():
        return redirect(url_for("dashboard"))

    @app.route("/praesentation")
    def praesentation():
        return render_template("praesentation.html")

    @app.route("/praesentation/kunde")
    def praesentation_kunde():
        return render_template("praesentation_kunde.html")

    @app.route("/praesentation/download")
    def praesentation_download():
        try:
            from pptx import Presentation as PPTXPresentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            abort(500, description="python-pptx nicht installiert.")

        def rgb(hex_str):
            hex_str = hex_str.lstrip('#')
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

        prs = PPTXPresentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_slide():
            return prs.slides.add_slide(blank)

        def rect(slide, x, y, w, h, fill=None):
            from pptx.util import Inches
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.line.fill.background()
            if fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(fill)
            else:
                shape.fill.background()
            return shape

        def txt(slide, text, x, y, w, h, size=24, bold=False, color='#000000', align='left', wrap=True):
            txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            txb.word_wrap = wrap
            tf = txb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
            return txb

        WS_DEEP   = '#250F6B'
        WS_BLUE   = '#451DC7'
        WS_GREEN  = '#04F06A'
        WHITE     = '#FFFFFF'
        GRAY      = '#595959'
        LIGHT_BG  = '#F0F2F5'

        # ── Folie 1: Cover ──────────────────────────────────────────
        s1 = add_slide()
        rect(s1, 0, 0, 13.33, 7.5, WS_DEEP)
        rect(s1, 0, 5.8, 13.33, 1.7, '#1a0a50')
        rect(s1, 0.5, 3.0, 0.08, 2.5, WS_GREEN)
        txt(s1, 'WARP', 0.8, 1.2, 10, 1.8, size=72, bold=True, color=WHITE)
        txt(s1, 'Wavestone Assessment Risk Planner', 0.8, 3.0, 10, 0.7, size=22, bold=False, color='#c5b8f5')
        txt(s1, 'Ihr Partner für strukturierte\nTest-Reifegrad-Bewertung', 0.8, 3.9, 9, 1.2, size=18, color='#ddd8f8')
        txt(s1, 'wavestone', 0.8, 6.5, 4, 0.5, size=14, bold=True, color=WS_GREEN)

        # ── Folie 2: WARP auf einen Blick (Was ist WARP + Ergebnisse) ──
        s2 = add_slide()
        rect(s2, 0, 0, 13.33, 7.5, LIGHT_BG)
        rect(s2, 0, 0, 13.33, 1.1, WS_BLUE)
        txt(s2, 'WARP auf einen Blick', 0.6, 0.18, 12, 0.75, size=28, bold=True, color=WHITE)
        # Left column: Was ist WARP?
        col_title_kw = {'size': 11, 'bold': True, 'color': WS_BLUE}
        txt(s2, 'WAS IST WARP?', 0.5, 1.2, 5.8, 0.4, **col_title_kw)
        ov_points = [
            ('Digitales Assessment-Framework', 'Strukturierte Bewertung des Test-Reifegrads auf Basis bewährter Standards — vollständig digital und papierlos.'),
            ('Schnell & standardisiert', 'Innerhalb von 2–3 Tagen zum vollständigen Überblick — reproduzierbar und projektübergreifend vergleichbar.'),
            ('Klare Ergebnisse', 'Konkreter Score, Stärken-/Schwächenprofil und direkt verwertbare Handlungsempfehlungen.'),
        ]
        for i, (title, desc) in enumerate(ov_points):
            by = 1.7 + i * 1.7
            rect(s2, 0.5, by, 5.8, 1.5, WHITE)
            txt(s2, title, 0.7, by + 0.15, 5.4, 0.45, size=13, bold=True, color=WS_DEEP)
            txt(s2, desc,  0.7, by + 0.65, 5.4, 0.75, size=11, color=GRAY)
        # Right column: Ihre Ergebnisse
        txt(s2, 'IHRE ERGEBNISSE', 7.0, 1.2, 5.8, 0.4, **col_title_kw)
        rect(s2, 7.0, 1.7, 5.8, 2.45, WHITE)
        rect(s2, 7.0, 1.7, 0.06, 2.45, WS_GREEN)
        txt(s2, 'PDF-Report', 7.2, 1.85, 5.4, 0.5, size=14, bold=True, color=WS_DEEP)
        report_items = ['Gesamtscore & Score je Testkategorie', 'Stärken- & Schwächenprofil', 'Konkrete Handlungsempfehlungen']
        for j, item in enumerate(report_items):
            txt(s2, '✓  ' + item, 7.2, 2.45 + j * 0.55, 5.4, 0.5, size=11, color=GRAY)
        rect(s2, 7.0, 4.3, 5.8, 2.55, WHITE)
        rect(s2, 7.0, 4.3, 0.06, 2.55, WS_GREEN)
        txt(s2, 'KI-generierte Testdokumente', 7.2, 4.45, 5.4, 0.5, size=14, bold=True, color=WS_DEEP)
        ki_docs = ['Teststrategie', 'Mastertestkonzept', 'Stufentestkonzept']
        for j, doc in enumerate(ki_docs):
            txt(s2, '✓  ' + doc, 7.2, 5.1 + j * 0.55, 5.4, 0.5, size=11, color=GRAY)

        # ── Folie 3: Nächste Schritte (CTA) ─────────────────────────
        s3 = add_slide()
        rect(s3, 0, 0, 13.33, 7.5, WS_DEEP)
        rect(s3, 0, 0, 13.33, 0.08, WS_GREEN)
        txt(s3, 'Starten Sie Ihr WARP-Assessment', 0.8, 0.5, 11.5, 1.0, size=30, bold=True, color=WHITE, align='center')
        steps2 = [('1', 'Termin vereinbaren', 'Gemeinsamen Workshop-Termin mit Ihrem Wavestone-Ansprechpartner festlegen.'),
                  ('2', 'Workshop durchführen', 'Beantwortung des Fragenkatalogs gemeinsam im Team — ca. 2–3 Tage.'),
                  ('3', 'Report & Roadmap', 'Vollständiger Report mit KI-Dokumenten und Empfehlungen als Ihre Roadmap.')]
        for i, (num, title, desc) in enumerate(steps2):
            bx = 0.7 + i * 4.0
            rect(s3, bx, 1.8, 3.6, 4.0, '#1a0a50')
            rect(s3, bx, 1.8, 3.6, 0.08, WS_GREEN)
            txt(s3, num,   bx + 1.55, 2.0,  0.5, 0.7, size=28, bold=True, color=WS_GREEN, align='center')
            txt(s3, title, bx + 0.2,  2.85, 3.2, 0.6, size=14, bold=True, color=WHITE, align='center')
            txt(s3, desc,  bx + 0.2,  3.6,  3.2, 1.9, size=11, color='#c5b8f5', align='center')
        txt(s3, 'wavestone', 0.8, 6.8, 12, 0.5, size=14, bold=True, color=WS_GREEN, align='center')

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(
            buf,
            as_attachment=True,
            download_name='WARP_Kundenpraesentation.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        )

    @app.route("/dashboard")
    @login_required
    def dashboard():
        projects = current_user.projects
        total_questions = _db_question_count()

        def project_progress(p):
            answered = sum(1 for a in p.answers if a.answer_id)
            return answered, total_questions

        open_projects = []
        done_projects = []
        for p in projects:
            answered, total = project_progress(p)
            pct = round(answered / total * 100) if total else 0
            entry = {"project": p, "answered": answered, "total": total, "pct": pct}
            (done_projects if p.is_complete else open_projects).append(entry)

        # Superuser: Gesamtübersicht aller User + Projekte
        all_users_data = None
        if current_user.is_superuser:
            all_users = db.session.execute(
                db.select(User).order_by(User.display_name, User.username)
            ).scalars().all()
            all_users_data = []
            for u in all_users:
                u_projects = []
                for p in u.projects:
                    answered, total = project_progress(p)
                    pct = round(answered / total * 100) if total else 0
                    u_projects.append({"project": p, "answered": answered, "total": total, "pct": pct})
                all_users_data.append({
                "user": u,
                "projects": u_projects,
                "open_count": sum(1 for e in u_projects if not e["project"].is_complete),
                "done_count": sum(1 for e in u_projects if e["project"].is_complete),
            })

        # Admin: Postkorb-Stats
        inbox_stats = None
        if current_user.is_admin:
            messages = db.session.execute(db.select(InboxMessage)).scalars().all()
            inbox_stats = {
                "neu": sum(1 for m in messages if m.status == "neu"),
                "in_bearbeitung": sum(1 for m in messages if m.status == "in_bearbeitung"),
                "erledigt": sum(1 for m in messages if m.status == "erledigt"),
            }

        neu_count = inbox_stats["neu"] if inbox_stats else 0
        return render_template(
            "dashboard.html",
            projects=projects,
            open_projects=open_projects,
            done_projects=done_projects,
            inbox_stats=inbox_stats,
            all_users_data=all_users_data,
            neu_count=neu_count,
        )

    @app.route("/project/<int:pid>/toggle-complete", methods=["POST"])
    @login_required
    def toggle_project_complete(pid: int):
        project = _get_project_or_403(pid)
        project.is_complete = not project.is_complete
        db.session.commit()
        return redirect(request.referrer or url_for("dashboard"))

    @app.route("/admin")
    @login_required
    def admin_overview():
        if not current_user.is_superuser:
            abort(403)
        users = db.session.execute(
            db.select(User).where(User.id != current_user.id).order_by(User.id)
        ).scalars().all()
        return render_template("admin.html", users=users, total_questions=_db_question_count())

    @app.route("/admin/user/new", methods=["POST"])
    @login_required
    def admin_user_new():
        if not current_user.is_superuser:
            abort(403)
        username = request.form.get("username", "").strip()
        display_name = request.form.get("display_name", "").strip()
        password = request.form.get("password", "").strip()
        role = request.form.get("role", "user")
        if role not in ('user', 'admin', 'superuser'):
            role = 'user'
        error = None
        if not username or not password:
            error = "Benutzername und Passwort sind erforderlich."
        elif len(password) < 6:
            error = "Passwort muss mindestens 6 Zeichen haben."
        else:
            existing = db.session.execute(
                db.select(User).where(User.username == username)
            ).scalar_one_or_none()
            if existing:
                error = f"Benutzername '{username}' ist bereits vergeben."
        if error:
            users = db.session.execute(
                db.select(User).where(User.id != current_user.id).order_by(User.id)
            ).scalars().all()
            return render_template("admin.html", users=users,
                                   total_questions=_db_question_count(), error=error)
        u = User(username=username, display_name=display_name or None, role=role)
        u.set_password(password)
        db.session.add(u)
        db.session.commit()
        return redirect(url_for("admin_overview"))

    @app.route("/admin/user/<int:uid>/set-role", methods=["POST"])
    @login_required
    def admin_user_set_role(uid: int):
        if not current_user.is_superuser:
            abort(403)
        if uid == current_user.id:
            abort(400)
        user = db.session.get(User, uid)
        if not user:
            abort(404)
        new_role = request.form.get("role", "user")
        if new_role not in ('user', 'admin', 'superuser'):
            abort(400)
        user.role = new_role
        db.session.commit()
        return redirect(url_for("admin_overview"))

    @app.route("/admin/user/<int:uid>/delete", methods=["POST"])
    @login_required
    def admin_user_delete(uid: int):
        if not current_user.is_superuser:
            abort(403)
        if uid == current_user.id:
            abort(400)
        user = db.session.get(User, uid)
        if not user:
            abort(404)
        project_ids = [p.id for p in user.projects]
        if project_ids:
            db.session.execute(
                db.update(InboxMessage)
                .where(InboxMessage.project_id.in_(project_ids))
                .values(project_id=None)
            )
            db.session.execute(
                db.delete(GeneratedDocument)
                .where(GeneratedDocument.project_id.in_(project_ids))
            )
        # Also null out InboxMessage.claimed_by_id references
        db.session.execute(
            db.update(InboxMessage)
            .where(InboxMessage.claimed_by_id == uid)
            .values(claimed_by_id=None)
        )
        db.session.delete(user)
        db.session.commit()
        return redirect(url_for("admin_overview"))

    @app.route("/admin/user/<int:uid>/toggle-lock", methods=["POST"])
    @login_required
    def admin_user_toggle_lock(uid: int):
        if not current_user.is_superuser:
            abort(403)
        if uid == current_user.id:
            abort(400)
        user = db.session.get(User, uid)
        if not user:
            abort(404)
        user.is_locked = not user.is_locked
        db.session.commit()
        return redirect(url_for("admin_overview"))

    # ------------------------------------------------------------------
    # Admin – Fragenkatalog
    # ------------------------------------------------------------------

    @app.route("/admin/questions")
    @login_required
    def admin_questions():
        if not current_user.is_superuser:
            abort(403)
        tab = request.args.get("tab", "assessment")
        if tab not in ("assessment", "workshop"):
            tab = "assessment"
        cats = db.session.execute(
            db.select(Category)
            .where(Category.catalog_type == tab)
            .order_by(Category.sort_order)
        ).scalars().all()
        return render_template(
            "admin_questions.html",
            categories=cats,
            levels=WARP_LEVEL_ORDER,
            total_questions=_db_question_count(tab),
            active_tab=tab,
        )

    @app.route("/admin/question/new", methods=["POST"])
    @login_required
    def admin_question_new():
        if not current_user.is_superuser:
            abort(403)
        category_id = request.form.get("category_id", "").strip()
        text = request.form.get("text", "").strip()
        hint = request.form.get("hint", "").strip() or None
        if not category_id or not text:
            flash("Fragetext darf nicht leer sein.", "error")
            return redirect(url_for("admin_questions"))
        cat = db.session.get(Category, category_id)
        if not cat:
            abort(404)
        max_sort = db.session.execute(
            db.select(db.func.max(Question.sort_order)).where(
                Question.category_id == category_id
            )
        ).scalar()
        new_sort = (max_sort or 0) + 1
        new_id = f"{category_id}-{uuid.uuid4().hex[:6]}"
        q = Question(
            id=new_id,
            category_id=category_id,
            text=text,
            hint=hint,
            sort_order=new_sort,
        )
        db.session.add(q)
        db.session.commit()
        flash("Frage erfolgreich hinzugefügt.", "success")
        return redirect(url_for("admin_questions") + f"#{category_id}")

    @app.route("/admin/question/<qid>/edit", methods=["POST"])
    @login_required
    def admin_question_edit(qid: str):
        if not current_user.is_superuser:
            abort(403)
        q = db.session.get(Question, qid)
        if not q:
            abort(404)
        new_text = request.form.get("text", "").strip()
        if new_text:
            q.text = new_text
        q.hint = request.form.get("hint", "").strip() or None
        db.session.commit()
        flash("Frage gespeichert.", "success")
        return redirect(url_for("admin_questions") + f"#{q.category_id}")

    @app.route("/admin/question/<qid>/delete", methods=["POST"])
    @login_required
    def admin_question_delete(qid: str):
        if not current_user.is_superuser:
            abort(403)
        q = db.session.get(Question, qid)
        if q:
            cat_id = q.category_id
            db.session.delete(q)
            db.session.commit()
            flash("Frage gelöscht.", "success")
            return redirect(url_for("admin_questions") + f"#{cat_id}")
        return redirect(url_for("admin_questions"))

    @app.route("/admin/category/new", methods=["POST"])
    @login_required
    def admin_category_new():
        if not current_user.is_superuser:
            abort(403)
        cat_id = request.form.get("id", "").strip()
        title = request.form.get("title", "").strip()
        parent = request.form.get("parent", "").strip()
        description = request.form.get("description", "").strip() or None
        if not cat_id or not title or not parent:
            flash("ID, Titel und Stufe/Bereich sind Pflichtfelder.", "error")
            return redirect(url_for("admin_questions"))
        if db.session.get(Category, cat_id):
            flash(f"Kategorie-ID '{cat_id}' ist bereits vergeben.", "error")
            return redirect(url_for("admin_questions"))
        catalog_type = request.form.get("catalog_type", "assessment")
        if catalog_type not in ("assessment", "workshop"):
            catalog_type = "assessment"
        if catalog_type == "assessment" and parent not in WARP_LEVEL_ORDER:
            flash("Ungültige WARP-Stufe.", "error")
            return redirect(url_for("admin_questions") + "?tab=assessment")
        max_sort = db.session.execute(
            db.select(db.func.max(Category.sort_order))
            .where(Category.catalog_type == catalog_type)
        ).scalar()
        cat = Category(
            id=cat_id,
            title=title,
            parent=parent,
            description=description,
            sort_order=(max_sort or 0) + 1,
            catalog_type=catalog_type,
        )
        db.session.add(cat)
        db.session.commit()
        flash(f"Kategorie '{title}' erfolgreich angelegt.", "success")
        return redirect(url_for("admin_questions") + f"?tab={catalog_type}#{cat_id}")

    @app.route("/admin/category/<cid>/edit", methods=["POST"])
    @login_required
    def admin_category_edit(cid: str):
        if not current_user.is_superuser:
            abort(403)
        cat = db.session.get(Category, cid)
        if not cat:
            abort(404)
        new_title = request.form.get("title", "").strip()
        if new_title:
            cat.title = new_title
        cat.description = request.form.get("description", "").strip() or None
        db.session.commit()
        flash("Kategorie gespeichert.", "success")
        return redirect(url_for("admin_questions") + f"#{cid}")

    @app.route("/admin/category/<cid>/delete", methods=["POST"])
    @login_required
    def admin_category_delete(cid: str):
        if not current_user.is_superuser:
            abort(403)
        cat = db.session.get(Category, cid)
        if cat:
            db.session.delete(cat)
            db.session.commit()
            flash("Kategorie und alle zugehörigen Fragen gelöscht.", "success")
        return redirect(url_for("admin_questions"))

    @app.route("/admin/workshop/import", methods=["POST"])
    @login_required
    def workshop_import():
        if not current_user.is_superuser:
            abort(403)
        file = request.files.get("file")
        if not file or not file.filename:
            flash("Keine Datei ausgewählt.", "error")
            return redirect(url_for("admin_questions") + "?tab=workshop")
        import csv as _csv
        import io as _io
        rows = []
        fname = file.filename.lower()
        try:
            if fname.endswith(".xlsx") or fname.endswith(".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(file, read_only=True, data_only=True)
                ws = wb.active
                for row in list(ws.iter_rows(min_row=2, values_only=True)):
                    rows.append([str(c).strip() if c is not None else "" for c in row])
            else:
                content = file.stream.read().decode("utf-8-sig")
                reader = _csv.reader(_io.StringIO(content))
                next(reader, None)
                rows = [r for r in reader]
        except Exception as exc:
            flash(f"Datei konnte nicht gelesen werden: {exc}", "error")
            return redirect(url_for("admin_questions") + "?tab=workshop")

        cats_created: dict[str, str] = {}
        questions_created = 0
        for row in rows:
            if not row or len(row) < 2:
                continue
            cat_name = str(row[0]).strip()
            q_text = str(row[1]).strip()
            q_hint = str(row[2]).strip() if len(row) > 2 and row[2] else None
            if not cat_name or not q_text:
                continue
            if cat_name not in cats_created:
                max_sort = db.session.execute(
                    db.select(db.func.max(Category.sort_order))
                    .where(Category.catalog_type == "workshop")
                ).scalar() or 0
                cat_id = f"ws-{uuid.uuid4().hex[:8]}"
                cat = Category(
                    id=cat_id, title=cat_name, parent="Workshop",
                    sort_order=max_sort + 1, catalog_type="workshop",
                )
                db.session.add(cat)
                db.session.flush()
                cats_created[cat_name] = cat_id
            cat_id = cats_created[cat_name]
            max_sort_q = db.session.execute(
                db.select(db.func.max(Question.sort_order))
                .where(Question.category_id == cat_id)
            ).scalar() or 0
            q = Question(
                id=f"{cat_id}-{uuid.uuid4().hex[:6]}",
                category_id=cat_id, text=q_text, hint=q_hint,
                sort_order=max_sort_q + 1,
            )
            db.session.add(q)
            questions_created += 1
        db.session.commit()
        flash(f"{questions_created} Fragen in {len(cats_created)} Kategorien importiert.", "success")
        return redirect(url_for("admin_questions") + "?tab=workshop")

    # ------------------------------------------------------------------
    # Project Routes
    # ------------------------------------------------------------------

    @app.route("/project/new", methods=["GET", "POST"])
    @login_required
    def new_project():
        if request.method == "POST":
            name = request.form.get("name", "").strip() or "Neues Projekt"
            catalog_type = request.form.get("catalog_type", "assessment")
            if catalog_type not in ("assessment", "workshop"):
                catalog_type = "assessment"
            project = Project(
                user_id=current_user.id,
                name=name,
                date=dt.date.today().strftime("%Y-%m-%d"),
                catalog_type=catalog_type,
            )
            db.session.add(project)
            db.session.commit()
            return redirect(url_for("questionnaire", pid=project.id))
        return render_template("new_project.html", projects=current_user.projects)

    @app.route("/project/<int:pid>")
    @login_required
    def questionnaire(pid: int):
        project = _get_project_or_403(pid)
        projects = current_user.projects
        ct = project.catalog_type or 'assessment'
        generated_docs = {
            d.doc_type: d.generated_at
            for d in db.session.execute(
                db.select(GeneratedDocument).where(GeneratedDocument.project_id == pid)
            ).scalars().all()
        }
        if ct == 'workshop':
            wq_list = project.workshop_questions
            categories = [{
                'id': f'ws-proj-{project.id}',
                'title': 'Workshop-Fragen',
                'description': '',
                'questions': [{'id': q.id, 'text': q.text, 'hint': q.hint or ''} for q in wq_list],
            }] if wq_list else []
            total_q = len(wq_list)
            can_manage = current_user.is_superuser or project.user_id == current_user.id
            workshop_questions = wq_list
        else:
            categories = _load_categories_from_db(ct)
            total_q = _db_question_count(ct)
            can_manage = False
            workshop_questions = []
        return render_template(
            "index.html",
            project=project,
            projects=projects,
            saved_answers=project.answers_dict(),
            saved_notes=project.notes_dict(),
            categories=categories,
            answer_options=ANSWER_OPTIONS,
            total_questions=total_q,
            today=dt.date.today().strftime("%Y-%m-%d"),
            generated_docs=generated_docs,
            can_manage_questions=can_manage,
            workshop_questions=workshop_questions,
        )

    @app.route("/project/<int:pid>/workshop/question/add", methods=["POST"])
    @login_required
    def workshop_question_add(pid: int):
        project = _get_project_or_403(pid)
        if project.catalog_type != 'workshop':
            abort(400)
        if not (current_user.is_superuser or project.user_id == current_user.id):
            abort(403)
        text = request.form.get("text", "").strip()
        if not text:
            return redirect(url_for("questionnaire", pid=pid) + "#verwalten")
        hint = request.form.get("hint", "").strip() or None
        max_order = db.session.execute(
            db.select(db.func.max(WorkshopQuestion.sort_order))
            .where(WorkshopQuestion.project_id == pid)
        ).scalar() or 0
        qid = f"wq-{pid}-{uuid.uuid4().hex[:8]}"
        q = WorkshopQuestion(id=qid, project_id=pid, text=text, hint=hint, sort_order=max_order + 1)
        db.session.add(q)
        db.session.commit()
        return redirect(url_for("questionnaire", pid=pid) + "#verwalten")

    @app.route("/project/<int:pid>/workshop/question/<qid>/edit", methods=["POST"])
    @login_required
    def workshop_question_edit(pid: int, qid: str):
        project = _get_project_or_403(pid)
        if not (current_user.is_superuser or project.user_id == current_user.id):
            abort(403)
        q = db.session.execute(
            db.select(WorkshopQuestion)
            .where(WorkshopQuestion.id == qid, WorkshopQuestion.project_id == pid)
        ).scalar_one_or_none()
        if q:
            text = request.form.get("text", "").strip()
            if text:
                q.text = text
            q.hint = request.form.get("hint", "").strip() or None
            db.session.commit()
        return redirect(url_for("questionnaire", pid=pid) + "#verwalten")

    @app.route("/project/<int:pid>/workshop/question/<qid>/delete", methods=["POST"])
    @login_required
    def workshop_question_delete(pid: int, qid: str):
        project = _get_project_or_403(pid)
        if not (current_user.is_superuser or project.user_id == current_user.id):
            abort(403)
        q = db.session.execute(
            db.select(WorkshopQuestion)
            .where(WorkshopQuestion.id == qid, WorkshopQuestion.project_id == pid)
        ).scalar_one_or_none()
        if q:
            db.session.execute(
                db.delete(Answer)
                .where(Answer.project_id == pid, Answer.question_id == qid)
            )
            db.session.delete(q)
            db.session.commit()
        return redirect(url_for("questionnaire", pid=pid) + "#verwalten")

    @app.route("/project/<int:pid>/answer", methods=["POST"])
    @login_required
    def save_answer(pid: int):
        project = _get_project_or_403(pid)
        data = request.get_json(silent=True) or {}
        question_id = data.get("question_id", "")
        answer_id = data.get("answer_id")
        note = data.get("note")
        if not question_id:
            return jsonify({"ok": False}), 400
        _upsert_answer(project.id, question_id, answer_id, note)
        return jsonify({"ok": True})

    @app.route("/project/<int:pid>/autofill", methods=["POST"])
    @login_required
    def autofill_project(pid: int):
        """Testhelfer: realistisches Zufallsszenario statt Gleichverteilung.

        Simuliert ein typisches Assessment: Stufe 2 wird zufällig entweder
        knapp erreicht oder knapp verfehlt (Basic), Stufe 3 ist nur ansatzweise
        erfüllt, und Stufe 4/5 haben kaum etwas vorzuweisen — statt gleichverteiltem
        Rauschen über alle vier Antwortoptionen.
        """
        import random
        project = _get_project_or_403(pid)

        stufe2_erreicht = random.choice([True, False])
        weights_by_level = {
            "Stufe 2 - Managed": (
                {"voll": 55, "teil": 32, "kaum": 10, "nicht": 3}
                if stufe2_erreicht else
                {"voll": 18, "teil": 32, "kaum": 32, "nicht": 18}
            ),
            "Stufe 3 - Defined": {"voll": 8, "teil": 24, "kaum": 34, "nicht": 34},
            "Stufe 4 - Measured": {"voll": 2, "teil": 8, "kaum": 20, "nicht": 70},
            "Stufe 5 - Optimization": {"voll": 1, "teil": 5, "kaum": 14, "nicht": 80},
        }
        default_weights = {"voll": 25, "teil": 25, "kaum": 25, "nicht": 25}

        categories = db.session.execute(
            db.select(Category).where(Category.catalog_type == "assessment")
        ).scalars().all()

        for cat in categories:
            weights = weights_by_level.get(cat.parent, default_weights)
            option_ids = list(weights.keys())
            option_weights = list(weights.values())
            for q in cat.questions:
                answer_id = random.choices(option_ids, weights=option_weights, k=1)[0]
                _upsert_answer(project.id, q.id, answer_id, None)
        db.session.commit()
        return redirect(url_for("questionnaire", pid=pid))

    @app.route("/project/<int:pid>/info", methods=["POST"])
    @login_required
    def save_project_info(pid: int):
        project = _get_project_or_403(pid)
        data = request.get_json(silent=True) or {}
        if "name" in data:
            project.name = data["name"].strip() or "Neues Projekt"
        if "owner" in data:
            project.owner = data["owner"].strip() or None
        if "date" in data:
            project.date = data["date"] or None
        db.session.commit()
        return jsonify({"ok": True})

    @app.route("/project/<int:pid>/delete", methods=["POST"])
    @login_required
    def delete_project(pid: int):
        project = _get_project_or_403(pid)
        owner_id = project.user_id
        # Null out InboxMessage.project_id before delete (nullable FK, no cascade configured)
        db.session.execute(
            db.update(InboxMessage)
            .where(InboxMessage.project_id == pid)
            .values(project_id=None)
        )
        # Delete generated documents (no cascade on Project side)
        db.session.execute(
            db.delete(GeneratedDocument)
            .where(GeneratedDocument.project_id == pid)
        )
        db.session.delete(project)
        db.session.commit()
        if current_user.is_superuser and owner_id != current_user.id:
            return redirect(url_for("admin_overview"))
        remaining = current_user.projects
        if remaining:
            return redirect(url_for("questionnaire", pid=remaining[0].id))
        return redirect(url_for("new_project"))

    @app.route("/project/<int:pid>/report/html", methods=["POST"])
    @login_required
    def project_report_html(pid: int):
        project = _get_project_or_403(pid)
        customer = request.form.get("variant") == "customer"
        ctx = _build_report_context(project.to_form_dict(), customer=customer)
        return render_template("report.html", **ctx)

    @app.route("/project/<int:pid>/report", methods=["POST"])
    @login_required
    def project_report_pdf(pid: int):
        try:
            from weasyprint import HTML  # type: ignore
        except ImportError:
            abort(500, description="WeasyPrint ist nicht installiert.")

        project = _get_project_or_403(pid)
        customer = request.form.get("variant") == "customer"
        ctx = _build_report_context(project.to_form_dict(), customer=customer)
        html = render_template("report.html", **ctx)
        pdf_bytes = HTML(string=html, base_url=request.host_url).write_pdf()
        prefix = "WARP_Kundenreport" if customer else "WARP_Report"
        filename = f"{prefix}_{project.name.replace(' ', '_')}_{dt.date.today().isoformat()}.pdf"
        return send_file(
            io.BytesIO(pdf_bytes),
            mimetype="application/pdf",
            as_attachment=True,
            download_name=filename,
        )

    # ------------------------------------------------------------------
    # KI-Dokumentgenerierung
    # ------------------------------------------------------------------

    _DOC_TEMPLATES = {
        'teststrategie':     'WARP_Teststrategie.docx',
        'mastertestkonzept': 'WARP_Mastertestkonzept.docx',
        'stufentestkonzept': 'WARP_Stufentestkonzept.docx',
    }
    _DOC_LABELS = {
        'teststrategie':     'Teststrategie',
        'mastertestkonzept': 'Mastertestkonzept',
        'stufentestkonzept': 'Stufentestkonzept',
    }

    def _build_ai_document(project: Project, doc_type: str) -> tuple[io.BytesIO, str]:
        import json as _json2
        import anthropic
        from docx import Document as DocxDoc
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn as _qn

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY ist nicht gesetzt.")

        answers = project.answers_dict()
        notes   = project.notes_dict()
        db_cats = _load_categories_from_db()

        label_lookup = {opt["id"]: opt["label"] for opt in ANSWER_OPTIONS}
        cat_lines, strengths, improvements = [], [], []
        weak_questions: list[str] = []

        for cat in db_cats:
            scores = [score_lookup.get(answers.get(q['id'], ''), 0) for q in cat['questions']]
            answered = [s for q, s in zip(cat['questions'], scores) if answers.get(q['id'])]
            if not answered:
                continue
            pct = round(sum(scores) / (len(cat['questions']) * 100) * 100)
            cat_lines.append(f"  - {cat['title']} ({cat['parent']}): {pct}%")
            if pct >= 70:
                strengths.append(cat['title'])
            elif pct < 40:
                improvements.append(cat['title'])
            for q in cat['questions']:
                ans_id = answers.get(q['id'], '')
                score  = score_lookup.get(ans_id, 0)
                note   = notes.get(q['id'], '')
                if ans_id and score < 50:
                    line = f"    [{cat['title']}] {q['text'][:110]} → {label_lookup.get(ans_id, ans_id)}"
                    if note:
                        line += f" | Notiz: {note[:80]}"
                    weak_questions.append(line)

        # Kapitelstruktur aus der Vorlage lesen
        tpl_path = ROOT / "app" / "static" / "docs" / _DOC_TEMPLATES[doc_type]
        doc = DocxDoc(str(tpl_path))
        sections: list[tuple[str, str]] = []
        cur_head = cur_desc = None
        for para in doc.paragraphs:
            if para.style.name == 'Heading 1':
                if cur_head:
                    sections.append((cur_head, cur_desc or ''))
                cur_head, cur_desc = para.text.strip(), None
            elif para.style.name == 'Normal' and cur_head and cur_desc is None:
                cur_desc = para.text.strip()
        if cur_head:
            sections.append((cur_head, cur_desc or ''))

        context = (
            f"Projekt: {project.name}\n"
            f"Verantwortliche/r: {project.owner or 'nicht angegeben'}\n"
            f"Datum: {project.date or 'nicht angegeben'}\n\n"
            f"WARP-Kategorie-Scores (Erfüllungsgrad je Prozessbereich):\n"
            + ('\n'.join(cat_lines) or '  (keine Antworten)')
            + f"\n\nStärken (≥70 %): {', '.join(strengths) or '—'}\n"
            f"Handlungsfelder (<40 %): {', '.join(improvements) or '—'}"
        )
        weak_ctx = (
            "\n\nEinzelfragen mit Handlungsbedarf (Score <50 %):\n"
            + '\n'.join(weak_questions[:40])
        ) if weak_questions else ""

        def _parse_chapter_response(raw: str) -> dict:
            raw = raw.strip()
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            return _json2.loads(raw.strip())

        # Ein API-Call pro Kapitel – verhindert JSON-Truncation bei langen Dokumenten
        client = anthropic.Anthropic(api_key=api_key)
        content: dict = {}
        for heading, desc in sections:
            chapter_prompt = f"""Du bist ein erfahrener Testmanagement-Berater bei Wavestone.
Erstelle auf Basis des folgenden WARP-Assessments professionellen deutschen Fachtext \
für das Kapitel „{heading}" eines „{_DOC_LABELS[doc_type]}"-Dokuments.

ASSESSMENT-ERGEBNISSE:
{context}{weak_ctx}

KAPITEL: {heading}
KAPITEL-BESCHREIBUNG: {desc}

ANFORDERUNGEN:
- Einleitender Fließtext (3–4 Sätze): bewertet den Ist-Zustand des Projekts in Bezug \
auf dieses Kapitel und begründet die Relevanz. Beziehe dich konkret auf die Score-Werte.
- 5–7 handlungsorientierte Aufzählungspunkte: jeder Punkt beschreibt eine konkrete \
Maßnahme oder Empfehlung mit WAS, WIE und WARUM (vollständiger Satz, ~30–50 Wörter).
- Verwende Fachbegriffe: risikobasiertes Testen, Äquivalenzklassen, Grenzwertanalyse, \
Shift-Left, Test Coverage, Defect-Escape-Rate, RACI, Regressionsstrategie usw.
- Tool-agnostisch: Tool-Kategorien beschreiben, Beispiele in Klammern erlaubt.
- Sprache: professionelles Deutsch, formeller Beratungsstil.
- Keine Zeilenumbrüche innerhalb von JSON-Strings. Keine einfachen oder doppelten \
Anführungszeichen innerhalb der Texte (stattdessen Langstriche oder Umschreibungen).

Antworte AUSSCHLIESSLICH mit diesem JSON, ohne Erklärungen:
{{"intro": "...", "bullets": ["...", "...", "...", "...", "..."]}}"""

            msg = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2048,
                messages=[{"role": "user", "content": chapter_prompt}],
            )
            try:
                content[heading] = _parse_chapter_response(msg.content[0].text)
            except Exception:
                content[heading] = {"intro": msg.content[0].text[:500], "bullets": []}

        # Hilfsfunktion: neuen Absatz direkt nach einem XML-Element einfügen
        def _insert_para_after(ref_p_xml, text: str, indent: bool = False):
            new_p = OxmlElement('w:p')
            if indent:
                pPr = OxmlElement('w:pPr')
                ind = OxmlElement('w:ind')
                ind.set(_qn('w:left'), '360')
                pPr.append(ind)
                new_p.append(pPr)
            r = OxmlElement('w:r')
            t = OxmlElement('w:t')
            t.text = text
            t.set(_qn('xml:space'), 'preserve')
            r.append(t)
            new_p.append(r)
            ref_p_xml.addnext(new_p)
            return new_p

        # Vorlage befüllen: Intro in Normal-Absatz, Bullets als eingefügte Absätze
        # list() einmalig auswerten – neu eingefügte Elemente sollen nicht mititeriert werden
        filled: set[str] = set()
        cur_head = None
        for para in list(doc.paragraphs):
            if para.style.name == 'Heading 1':
                cur_head = para.text.strip()
            elif para.style.name == 'Normal' and cur_head and cur_head not in filled:
                chapter_data = content.get(cur_head)
                if not chapter_data:
                    continue
                intro   = chapter_data.get('intro', '') if isinstance(chapter_data, dict) else str(chapter_data)
                bullets = chapter_data.get('bullets', []) if isinstance(chapter_data, dict) else []

                # Intro in bestehenden Normal-Absatz schreiben
                for run in para.runs:
                    run.text = ''
                if para.runs:
                    para.runs[0].text = intro
                else:
                    para.add_run(intro)

                # Bullets als neue Absätze nach dem Intro einfügen (umgekehrte Reihenfolge wegen addnext)
                last_xml = para._p
                for bullet in bullets:
                    new_xml = _insert_para_after(last_xml, f'•  {bullet}', indent=True)
                    last_xml = new_xml

                filled.add(cur_head)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        label  = _DOC_LABELS[doc_type]
        fname  = f"WARP_{label}_{project.name.replace(' ', '_')}.docx"
        return buf, fname

    @app.route("/project/<int:pid>/generate/<doc_type>", methods=["POST"])
    @login_required
    def generate_document(pid: int, doc_type: str):
        if doc_type not in _DOC_TEMPLATES:
            abort(404)
        project = _get_project_or_403(pid)
        try:
            buf, fname = _build_ai_document(project, doc_type)
            file_bytes = buf.getvalue()
            # Upsert: vorhandenes Dokument ersetzen oder neu anlegen
            existing = db.session.execute(
                db.select(GeneratedDocument).where(
                    GeneratedDocument.project_id == pid,
                    GeneratedDocument.doc_type == doc_type,
                )
            ).scalar_one_or_none()
            if existing:
                existing.file_data = file_bytes
                existing.filename = fname
                existing.generated_at = dt.datetime.utcnow()
            else:
                db.session.add(GeneratedDocument(
                    project_id=pid,
                    doc_type=doc_type,
                    filename=fname,
                    file_data=file_bytes,
                ))
            db.session.commit()
            return send_file(
                io.BytesIO(file_bytes),
                mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                as_attachment=True,
                download_name=fname,
            )
        except ValueError as exc:
            return render_template("error.html", message=str(exc)), 500
        except Exception as exc:
            import traceback
            print(traceback.format_exc())
            return render_template("error.html", message=f"Generierung fehlgeschlagen: {exc}"), 500

    @app.route("/project/<int:pid>/document/<doc_type>")
    @login_required
    def download_generated_document(pid: int, doc_type: str):
        _get_project_or_403(pid)
        doc = db.session.execute(
            db.select(GeneratedDocument).where(
                GeneratedDocument.project_id == pid,
                GeneratedDocument.doc_type == doc_type,
            )
        ).scalar_one_or_none()
        if not doc:
            abort(404)
        return send_file(
            io.BytesIO(doc.file_data),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=doc.filename,
        )

    # ------------------------------------------------------------------
    # Vorlagen-Seite
    # ------------------------------------------------------------------

    @app.route("/vorlagen")
    @login_required
    def vorlagen():
        assessment_projects = [p for p in current_user.projects if p.catalog_type == 'assessment']
        total_questions = _db_question_count('assessment')

        projects_data = []
        for p in assessment_projects:
            answered = sum(1 for a in p.answers if a.answer_id)
            pct = round(answered / total_questions * 100) if total_questions else 0
            gen_docs = {
                d.doc_type: d.generated_at.strftime('%d.%m.%Y')
                for d in db.session.execute(
                    db.select(GeneratedDocument).where(GeneratedDocument.project_id == p.id)
                ).scalars().all()
            }
            projects_data.append({
                'project': p,
                'answered': answered,
                'total': total_questions,
                'pct': pct,
                'is_100': total_questions > 0 and answered >= total_questions,
                'generated_docs': gen_docs,
            })

        neu_count = 0
        if current_user.is_admin:
            neu_count = db.session.execute(
                db.select(db.func.count()).select_from(InboxMessage).where(
                    InboxMessage.status == 'neu'
                )
            ).scalar() or 0

        return render_template(
            'vorlagen.html',
            projects=current_user.projects,
            projects_data=projects_data,
            neu_count=neu_count,
        )

    # ------------------------------------------------------------------
    # Inbox – Webhook + Admin-Postkorb
    # ------------------------------------------------------------------

    @app.route("/api/inbox", methods=["POST"])
    @csrf.exempt
    def api_inbox_receive():
        api_key = os.environ.get("WARP_INBOX_API_KEY", "")
        if api_key:
            provided = request.headers.get("X-API-Key", "")
            if not provided or provided != api_key:
                return jsonify({"error": "Unauthorized"}), 401
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "Invalid JSON"}), 400
        user_name = data.get("userName", "").strip()
        user_email = data.get("userEmail", "").strip()
        recommendation = data.get("recommendation", "").strip()
        if not user_name or not user_email or not recommendation:
            return jsonify({"error": "Missing required fields: userName, userEmail, recommendation"}), 400
        msg = InboxMessage(
            source=data.get("source", "IMPULSE"),
            user_name=user_name,
            user_email=user_email,
            recommendation=recommendation,
            scores_json=_json.dumps(data.get("scores", {})),
            rationale=data.get("rationale", ""),
            top_factors_json=_json.dumps(data.get("topFactors", [])),
            maturity=data.get("maturity"),
            contact_pref=data.get("contactPref"),
            contact_phone=data.get("contactPhone"),
        )
        db.session.add(msg)
        db.session.commit()
        print(f"[WARP Postkorb] Neue Anfrage von {user_name} ({user_email}): {recommendation}")
        return jsonify({"ok": True, "id": msg.id}), 201

    @app.route("/api/inbox/count")
    @login_required
    def api_inbox_count():
        if not current_user.is_admin:
            return jsonify({"count": 0})
        count = db.session.execute(
            db.select(db.func.count()).select_from(InboxMessage).where(
                InboxMessage.status == "neu"
            )
        ).scalar() or 0
        return jsonify({"count": count})

    @app.route("/inbox")
    @login_required
    def inbox():
        if not current_user.is_admin:
            abort(403)
        messages = db.session.execute(
            db.select(InboxMessage).order_by(InboxMessage.received_at.desc())
        ).scalars().all()
        for m in messages:
            m.scores = _json.loads(m.scores_json or "{}")
            m.top_factors = _json.loads(m.top_factors_json or "[]")
        neu_count = sum(1 for m in messages if m.status == "neu")
        all_users = db.session.execute(
            db.select(User).order_by(User.display_name, User.username)
        ).scalars().all()
        projects = current_user.projects
        first_pid = projects[0].id if projects else None
        generated_docs = {}
        if first_pid:
            generated_docs = {
                d.doc_type: d.generated_at
                for d in db.session.execute(
                    db.select(GeneratedDocument).where(GeneratedDocument.project_id == first_pid)
                ).scalars().all()
            }
        return render_template("inbox.html", messages=messages, neu_count=neu_count,
                               admins=all_users, projects=projects, first_pid=first_pid,
                               generated_docs=generated_docs)

    @app.route("/inbox/<int:mid>/assign", methods=["POST"])
    @login_required
    def inbox_assign(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if not msg:
            abort(404)
        try:
            assignee_id = int(request.form.get("assignee_id", 0))
        except ValueError:
            abort(400)
        assignee = db.session.get(User, assignee_id)
        if not assignee:
            abort(400)
        msg.status = "in_bearbeitung"
        msg.claimed_by_id = assignee_id
        msg.claimed_at = dt.datetime.utcnow()
        # Projekt mitumhängen falls vorhanden
        if msg.project_id:
            project = db.session.get(Project, msg.project_id)
            if project:
                project.user_id = assignee_id
        db.session.commit()
        return redirect(url_for("inbox"))

    @app.route("/inbox/<int:mid>/create-project", methods=["POST"])
    @login_required
    def inbox_create_project(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if not msg:
            abort(404)
        # Serverseitiger Guard: Projekt bereits vorhanden → zurück zum Postkorb
        if msg.project_id:
            return redirect(url_for("inbox"))
        # Projekt-Besitzer: zugewiesener User falls vorhanden, sonst aktueller User
        owner_user = db.session.get(User, msg.claimed_by_id) if msg.claimed_by_id else current_user
        project = Project(
            user_id=owner_user.id,
            name=f"Assessment – {msg.user_name}",
            owner=msg.user_name,
            date=dt.date.today().isoformat(),
        )
        db.session.add(project)
        db.session.commit()  # Commit first so project.id is available
        # Nachricht mit Projekt verknüpfen
        msg.project_id = project.id
        if msg.status == "neu":
            msg.status = "in_bearbeitung"
            msg.claimed_by_id = owner_user.id
            msg.claimed_at = dt.datetime.utcnow()
        db.session.commit()
        return redirect(url_for("inbox"))

    @app.route("/inbox/<int:mid>/claim", methods=["POST"])
    @login_required
    def inbox_claim(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if not msg:
            abort(404)
        if msg.status == "neu":
            msg.status = "in_bearbeitung"
            msg.claimed_by_id = current_user.id
            msg.claimed_at = dt.datetime.utcnow()
            db.session.commit()
        return redirect(url_for("inbox"))

    @app.route("/inbox/<int:mid>/done", methods=["POST"])
    @login_required
    def inbox_done(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if not msg:
            abort(404)
        msg.status = "erledigt"
        db.session.commit()
        return redirect(url_for("inbox"))

    @app.route("/inbox/<int:mid>/release", methods=["POST"])
    @login_required
    def inbox_release(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if not msg:
            abort(404)
        msg.status = "neu"
        msg.claimed_by_id = None
        msg.claimed_at = None
        db.session.commit()
        return redirect(url_for("inbox"))

    @app.route("/inbox/<int:mid>/delete", methods=["POST"])
    @login_required
    def inbox_delete(mid: int):
        if not current_user.is_admin:
            abort(403)
        msg = db.session.get(InboxMessage, mid)
        if msg:
            db.session.delete(msg)
            db.session.commit()
        return redirect(url_for("inbox"))

    # ------------------------------------------------------------------
    # Security Headers
    # ------------------------------------------------------------------

    @app.after_request
    def set_security_headers(response):
        response.headers["X-Frame-Options"] = "DENY"
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
        response.headers["Content-Security-Policy"] = (
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline'; "
            "style-src 'self' 'unsafe-inline' fonts.googleapis.com; "
            "font-src 'self' fonts.gstatic.com data:; "
            "img-src 'self' data:; "
            "connect-src 'self'; "
            "frame-ancestors 'none';"
        )
        return response

    # ------------------------------------------------------------------
    # Coverage Matrix
    # ------------------------------------------------------------------

    @app.route("/coverage-matrix")
    def coverage_matrix():
        from .data.questions import CATEGORIES as PY_CATS, ANSWER_OPTIONS

        # E2E-Testfälle (Playwright)
        e2e_tests = [
            # Login
            {"id": "TC-LOGIN-P1", "area": "Login", "type": "E2E", "desc": "Alle UI-Elemente der Login-Seite sind sichtbar"},
            {"id": "TC-LOGIN-P2", "area": "Login", "type": "E2E", "desc": "Admin-Login leitet auf /admin weiter"},
            {"id": "TC-LOGIN-P3", "area": "Login", "type": "E2E", "desc": "User-Login leitet auf /project/... weiter"},
            {"id": "TC-LOGIN-N1", "area": "Login", "type": "E2E", "desc": "Falsches Passwort → Fehlermeldung"},
            {"id": "TC-LOGIN-N2", "area": "Login", "type": "E2E", "desc": "Nicht existenter Benutzer → Fehlermeldung"},
            {"id": "TC-LOGIN-N3", "area": "Login", "type": "E2E", "desc": "Gesperrter Benutzer → spezifische Meldung"},
            # Register
            {"id": "TC-REG-P1", "area": "Registrierung", "type": "E2E", "desc": "Alle UI-Elemente der Registrierungsseite sichtbar"},
            {"id": "TC-REG-P2", "area": "Registrierung", "type": "E2E", "desc": "Neuer Benutzer wird angelegt und auf Projekt weitergeleitet"},
            {"id": "TC-REG-P3", "area": "Registrierung", "type": "E2E", "desc": "Link 'Bereits registriert?' navigiert zu /login"},
            {"id": "TC-REG-N1", "area": "Registrierung", "type": "E2E", "desc": "Doppelter Benutzername → Fehlermeldung"},
            {"id": "TC-REG-N2", "area": "Registrierung", "type": "E2E", "desc": "Passwörter stimmen nicht überein → Fehler"},
            {"id": "TC-REG-N3", "area": "Registrierung", "type": "E2E", "desc": "Passwort kürzer als 6 Zeichen → Fehler"},
            # Admin
            {"id": "TC-ADMIN-P1", "area": "Administration", "type": "E2E", "desc": "Alle UI-Elemente der Admin-Seite sichtbar"},
            {"id": "TC-ADMIN-P2", "area": "Administration", "type": "E2E", "desc": "Admin kann neuen Benutzer anlegen"},
            {"id": "TC-ADMIN-P3", "area": "Administration", "type": "E2E", "desc": "Admin kann Benutzer sperren – Badge erscheint"},
            {"id": "TC-ADMIN-N1", "area": "Administration", "type": "E2E", "desc": "Normaler Benutzer kann /admin nicht aufrufen (403)"},
            {"id": "TC-ADMIN-N2", "area": "Administration", "type": "E2E", "desc": "Nicht eingeloggter Benutzer wird zu /login umgeleitet"},
            {"id": "TC-ADMIN-N3", "area": "Administration", "type": "E2E", "desc": "Admin kann kein neues Projekt anlegen"},
            # Questionnaire
            {"id": "TC-QUEST-P1", "area": "Fragenkatalog", "type": "E2E", "desc": "Alle UI-Elemente des Fragenkatalogs sichtbar"},
            {"id": "TC-QUEST-P2", "area": "Fragenkatalog", "type": "E2E", "desc": "Antwort wird gespeichert (AJAX) und bleibt nach Reload"},
            {"id": "TC-QUEST-P3", "area": "Fragenkatalog", "type": "E2E", "desc": "Download-Buttons für Vorlagen vorhanden und verlinkt"},
            {"id": "TC-QUEST-N1", "area": "Fragenkatalog", "type": "E2E", "desc": "Nicht eingeloggter Benutzer wird zu /login umgeleitet"},
            {"id": "TC-QUEST-N2", "area": "Fragenkatalog", "type": "E2E", "desc": "Fremdes Projekt liefert 403"},
            {"id": "TC-QUEST-N3", "area": "Fragenkatalog", "type": "E2E", "desc": "Nicht existierende Projekt-ID liefert 403"},
        ]

        # Unit-Testfälle (Flask Test Client)
        unit_tests = [
            {"id": "TC-AUTH-01", "area": "Auth", "type": "Unit", "desc": "GET /login → HTTP 200"},
            {"id": "TC-AUTH-02", "area": "Auth", "type": "Unit", "desc": "GET /register → HTTP 200"},
            {"id": "TC-AUTH-03", "area": "Auth", "type": "Unit", "desc": "POST /login gültig → Redirect 302"},
            {"id": "TC-AUTH-04", "area": "Auth", "type": "Unit", "desc": "POST /login ungültig → Fehlermeldung"},
            {"id": "TC-AUTH-05", "area": "Auth", "type": "Unit", "desc": "GET /logout → Redirect"},
            {"id": "TC-AUTH-06", "area": "Auth", "type": "Unit", "desc": "GET /dashboard ohne Login → Redirect /login"},
            {"id": "TC-AUTH-07", "area": "Auth", "type": "Unit", "desc": "GET /admin ohne Login → Redirect /login"},
            {"id": "TC-AUTH-08", "area": "Auth", "type": "Unit", "desc": "GET /admin als normaler User → 403"},
            {"id": "TC-AUTH-09", "area": "Auth", "type": "Unit", "desc": "POST /register doppelter Username → Fehler"},
            {"id": "TC-AUTH-10", "area": "Auth", "type": "Unit", "desc": "POST /register Passwörter stimmen nicht überein"},
            {"id": "TC-AUTH-11", "area": "Auth", "type": "Unit", "desc": "POST /register Passwort zu kurz → Fehler"},
            {"id": "TC-PROJ-01", "area": "Projekte", "type": "Unit", "desc": "POST /project/new → Projekt angelegt, Redirect"},
            {"id": "TC-PROJ-02", "area": "Projekte", "type": "Unit", "desc": "GET /project/<id> → Fragenkatalog 200"},
            {"id": "TC-PROJ-03", "area": "Projekte", "type": "Unit", "desc": "POST /project/<id>/answer → Antwort gespeichert"},
            {"id": "TC-PROJ-04", "area": "Projekte", "type": "Unit", "desc": "POST /project/<id>/answer erneut → Upsert"},
            {"id": "TC-PROJ-05", "area": "Projekte", "type": "Unit", "desc": "POST /project/<id>/info → Infos aktualisiert"},
            {"id": "TC-PROJ-06", "area": "Projekte", "type": "Unit", "desc": "GET fremdes Projekt → 403"},
            {"id": "TC-PROJ-07", "area": "Projekte", "type": "Unit", "desc": "GET /project/999999 → 403"},
            {"id": "TC-PROJ-08", "area": "Projekte", "type": "Unit", "desc": "GET /dashboard eingeloggt → 200"},
            {"id": "TC-PROJ-09", "area": "Projekte", "type": "Unit", "desc": "POST /project/<id>/report/html → HTML-Report 200"},
            {"id": "TC-PROJ-10", "area": "Projekte", "type": "Unit", "desc": "GET /project/new → Formular 200"},
            {"id": "TC-PROJ-11", "area": "Projekte", "type": "Unit", "desc": "POST /project/<id>/complete → Status gewechselt"},
            {"id": "TC-API-01", "area": "API/Webhook", "type": "Unit", "desc": "POST /api/inbox gültiger Key → 201"},
            {"id": "TC-API-02", "area": "API/Webhook", "type": "Unit", "desc": "POST /api/inbox falscher Key → 401"},
            {"id": "TC-API-03", "area": "API/Webhook", "type": "Unit", "desc": "POST /api/inbox fehlendes Feld → 400"},
            {"id": "TC-API-04", "area": "API/Webhook", "type": "Unit", "desc": "GET /api/inbox/count Admin → JSON count"},
            {"id": "TC-API-05", "area": "API/Webhook", "type": "Unit", "desc": "GET /api/inbox/count anonym → 302"},
            {"id": "TC-API-06", "area": "Postkorb", "type": "Unit", "desc": "GET /inbox Admin → 200"},
            {"id": "TC-API-07", "area": "Postkorb", "type": "Unit", "desc": "GET /inbox normaler User → 403"},
            {"id": "TC-API-08", "area": "Postkorb", "type": "Unit", "desc": "GET /admin/questions Superuser → 200"},
            {"id": "TC-API-09", "area": "Postkorb",  "type": "Unit", "desc": "GET /coverage-matrix oeffentlich → 200"},
            {"id": "TC-SEC-01", "area": "Security",  "type": "Unit", "desc": "SQL-Injection im Login → kein Auth-Bypass"},
            {"id": "TC-SEC-02", "area": "Security",  "type": "Unit", "desc": "XSS im Projektnamen → HTML-escaped"},
            {"id": "TC-SEC-03", "area": "Security",  "type": "Unit", "desc": "CSRF-Extension registriert + POST ohne Token → 400"},
            {"id": "TC-SEC-04", "area": "Security",  "type": "Unit", "desc": "Flask-Login initialisiert + session_protection aktiv"},
            {"id": "TC-SEC-05", "area": "Security",  "type": "Unit", "desc": "Path-Traversal via Static-URL → blockiert"},
            {"id": "TC-SEC-06", "area": "Security",  "type": "Unit", "desc": "10x Fehlanmeldung → kein HTTP 500"},
            {"id": "TC-SEC-07", "area": "Security",  "type": "Unit", "desc": "Cross-User Projektzugriff → 403"},
            {"id": "TC-SEC-08", "area": "Security",  "type": "Unit", "desc": "Sicherheits-Header vorhanden (X-Frame-Options etc.)"},
            {"id": "TC-SEC-09", "area": "Security",  "type": "Unit", "desc": "API-Key nicht in Response-Body sichtbar"},
            {"id": "TC-SEC-10", "area": "Security",  "type": "Unit", "desc": "Alle geschuetzten Routen → Redirect /login"},
            {"id": "TC-SEC-11", "area": "Security",  "type": "Unit", "desc": "Admin-Routen fuer normalen User → 403"},
        ]

        # Fragenkatalog-Statistik
        catalog = {}
        for cat in PY_CATS:
            level = cat["parent"]
            if level not in catalog:
                catalog[level] = []
            catalog[level].append({
                "id": cat["id"],
                "title": cat["title"],
                "total": len(cat["questions"]),
                "new": sum(1 for q in cat["questions"] if q.get("new")),
            })

        total_questions = sum(len(c["questions"]) for c in PY_CATS)
        total_e2e = len(e2e_tests)
        total_unit = len(unit_tests)

        return render_template(
            "coverage_matrix.html",
            e2e_tests=e2e_tests,
            unit_tests=unit_tests,
            catalog=catalog,
            total_questions=total_questions,
            total_e2e=total_e2e,
            total_unit=total_unit,
        )

    return app


app = create_app()


if __name__ == "__main__":
    app.run(debug=True, port=5000)
